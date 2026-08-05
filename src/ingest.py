import os
from os import listdir
from os.path import isfile, join, isdir

os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")
os.environ.setdefault("OMP_NUM_THREADS", "1")

import torch
torch.set_num_threads(1)

import chromadb
import docx
from pptx import Presentation
from pypdf import PdfReader
from langchain_text_splitters import TokenTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_chroma import Chroma

DB_PATH = "./chroma_db"
DOCS_PATH = "./docs"
COLLECTION_NAME = "treinamento_rag"
EMBEDDING_MODEL = "paraphrase-multilingual-MiniLM-L12-v2"

def lista_arquivos(dir):
    arquivo_list = []
    for f in listdir(dir):
        if isfile(join(dir, f)):
            arquivo_list.append(join(dir, f))
        elif isdir(join(dir, f)):
            arquivo_list += lista_arquivos(join(dir, f))
    return arquivo_list

def carrega_texto_word(path):
    doc = docx.Document(path)
    return "\n".join([para.text for para in doc.paragraphs])

def carrega_texto_pptx(path):
    prs = Presentation(path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return "\n".join(full_text)

def main_indexing():
    os.makedirs(DOCS_PATH, exist_ok=True)

    print("🧠 Carregando modelo de embeddings paraphrase-multilingual-MiniLM-L12-v2 (384 dims)...")
    hf = HuggingFaceEmbeddings(
        model_name=EMBEDDING_MODEL,
        model_kwargs={"device": "cpu"},
        encode_kwargs={"normalize_embeddings": True},
    )

    print("🗑️  Limpando coleções antigas do ChromaDB...")
    client = chromadb.PersistentClient(path=DB_PATH)
    for c in client.list_collections():
        client.delete_collection(c.name)

    vectorstore = Chroma(
        collection_name=COLLECTION_NAME,
        embedding_function=hf,
        persist_directory=DB_PATH,
        collection_metadata={"hnsw:space": "cosine"},
    )

    print(f"\n📂 Indexando os documentos de {DOCS_PATH}...\n")
    arquivos = lista_arquivos(DOCS_PATH)

    total_chunks = 0
    text_splitter = TokenTextSplitter(chunk_size=250, chunk_overlap=50)

    for arquivo in arquivos:
        try:
            arquivo_content = ""

            if arquivo.lower().endswith(".pdf"):
                print("Indexando: " + arquivo)
                reader = PdfReader(arquivo)
                for page in reader.pages:
                    arquivo_content += " " + (page.extract_text() or "")

            elif arquivo.lower().endswith((".txt", ".md")):
                print("Indexando: " + arquivo)
                with open(arquivo, "r", encoding="utf-8") as f:
                    arquivo_content = f.read()

            elif arquivo.lower().endswith(".docx"):
                print("Indexando: " + arquivo)
                arquivo_content = carrega_texto_word(arquivo)

            elif arquivo.lower().endswith(".pptx"):
                print("Indexando: " + arquivo)
                arquivo_content = carrega_texto_pptx(arquivo)

            else:
                continue

            if not arquivo_content.strip():
                print("  ⚠️  Nenhum texto extraído (possível PDF escaneado). Pulando.")
                continue

            textos = text_splitter.split_text(arquivo_content)
            metadata = [{"path": arquivo} for _ in textos]
            vectorstore.add_texts(textos, metadatas=metadata)
            total_chunks += len(textos)
            print(f"  -> {len(textos)} chunks adicionados")

        except Exception as e:
            print(f"O processo falhou para o arquivo {arquivo}: {e}")

    print(f"\n✅ Indexação concluída! {total_chunks} chunks no total.")

if __name__ == "__main__":
    main_indexing()
